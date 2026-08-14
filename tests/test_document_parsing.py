"""Tests for document text extraction utilities."""

from pathlib import Path

from durin.utils.document import (
    SUPPORTED_EXTENSIONS,
    _is_text_extension,
    extract_text,
)


class TestSupportedExtensions:
    """Test the SUPPORTED_EXTENSIONS constant."""

    def test_supported_extensions_include_common_formats(self):
        """Test that common document formats are included."""
        # Document formats
        assert ".pdf" in SUPPORTED_EXTENSIONS
        assert ".docx" in SUPPORTED_EXTENSIONS
        assert ".xlsx" in SUPPORTED_EXTENSIONS
        assert ".pptx" in SUPPORTED_EXTENSIONS

        # Text formats
        assert ".txt" in SUPPORTED_EXTENSIONS
        assert ".md" in SUPPORTED_EXTENSIONS
        assert ".csv" in SUPPORTED_EXTENSIONS
        assert ".json" in SUPPORTED_EXTENSIONS
        assert ".yaml" in SUPPORTED_EXTENSIONS
        assert ".yml" in SUPPORTED_EXTENSIONS

        # Image formats
        assert ".png" in SUPPORTED_EXTENSIONS
        assert ".jpg" in SUPPORTED_EXTENSIONS
        assert ".jpeg" in SUPPORTED_EXTENSIONS


class TestExtractText:
    """Test the extract_text function."""

    def test_extract_text_unsupported_returns_none(self, tmp_path: Path):
        """Test that unsupported file types return None."""
        unsupported_file = tmp_path / "file.xyz"
        unsupported_file.write_text("content")

        result = extract_text(unsupported_file)
        assert result is None

    def test_extract_text_file_not_found(self, tmp_path: Path):
        """Test that non-existent files return error string."""
        missing_file = tmp_path / "nonexistent.txt"

        result = extract_text(missing_file)
        assert result is not None
        assert "[error: file not found:" in result

    def test_extract_text_txt_file(self, tmp_path: Path):
        """Test extracting text from a .txt file."""
        txt_file = tmp_path / "test.txt"
        content = "Hello, world!\nThis is a test."
        txt_file.write_text(content, encoding="utf-8")

        result = extract_text(txt_file)
        assert result == content

    def test_extract_text_txt_file_with_truncation(self, tmp_path: Path):
        """Test that large text files are truncated."""
        txt_file = tmp_path / "large.txt"
        # Create content larger than _MAX_TEXT_LENGTH
        content = "x" * 300_000
        txt_file.write_text(content, encoding="utf-8")

        result = extract_text(txt_file)
        assert len(result) < 300_000
        assert "(truncated," in result
        assert "chars total)" in result

    def test_extract_text_md_file(self, tmp_path: Path):
        """Test extracting text from a .md file."""
        md_file = tmp_path / "test.md"
        content = "# Header\n\nSome markdown content."
        md_file.write_text(content, encoding="utf-8")

        result = extract_text(md_file)
        assert result == content

    def test_extract_text_csv_file(self, tmp_path: Path):
        """Test extracting text from a .csv file."""
        csv_file = tmp_path / "test.csv"
        content = "name,age\nAlice,30\nBob,25"
        csv_file.write_text(content, encoding="utf-8")

        result = extract_text(csv_file)
        assert result == content

    def test_extract_text_json_file(self, tmp_path: Path):
        """Test extracting text from a .json file."""
        json_file = tmp_path / "test.json"
        content = '{"key": "value", "number": 42}'
        json_file.write_text(content, encoding="utf-8")

        result = extract_text(json_file)
        assert result == content

    def test_extract_text_xlsx(self, tmp_path: Path):
        """Test extracting text from an .xlsx file."""
        from openpyxl import Workbook

        xlsx_file = tmp_path / "test.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Name"
        ws["B1"] = "Age"
        ws["A2"] = "Alice"
        ws["B2"] = 30
        ws["A3"] = "Bob"
        ws["B3"] = 25

        # Add a second sheet
        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = "Product"
        ws2["B1"] = "Price"
        ws2["A2"] = "Widget"
        ws2["B2"] = 9.99

        wb.save(xlsx_file)
        wb.close()

        result = extract_text(xlsx_file)
        assert result is not None
        # markitdown renders each sheet as a "## <name>" heading + a table.
        assert "Sheet1" in result
        assert "Sheet2" in result
        assert "Alice" in result
        assert "Bob" in result
        assert "Widget" in result
        assert "9.99" in result

    def test_extract_text_xlsx_empty_sheet(self, tmp_path: Path):
        """Test extracting text from an .xlsx file with empty sheets."""
        from openpyxl import Workbook

        xlsx_file = tmp_path / "empty.xlsx"
        wb = Workbook()
        # Clear the default sheet
        wb.remove(wb.active)
        # Add an empty sheet
        wb.create_sheet("EmptySheet")
        wb.save(xlsx_file)
        wb.close()

        result = extract_text(xlsx_file)
        # An empty sheet yields the sheet heading with an empty table, or an
        # error / empty string — the point is no data content leaks.
        assert (
            "EmptySheet" in result
            or result.startswith("[error:")
            or result == ""
        )

    def test_extract_text_docx(self, tmp_path: Path):
        """Test extracting text from a .docx file."""
        from docx import Document

        docx_file = tmp_path / "test.docx"
        doc = Document()
        doc.add_heading("Test Document", 0)
        doc.add_paragraph("This is paragraph one.")
        doc.add_paragraph("This is paragraph two.")
        doc.save(docx_file)

        result = extract_text(docx_file)
        assert result is not None
        assert "Test Document" in result
        assert "This is paragraph one." in result
        assert "This is paragraph two." in result

    def test_extract_text_docx_empty(self, tmp_path: Path):
        """Test extracting text from an empty .docx file."""
        from docx import Document

        docx_file = tmp_path / "empty.docx"
        doc = Document()
        doc.save(docx_file)

        result = extract_text(docx_file)
        # An empty document has no extractable text — markitdown reports it as
        # an error string (filtered downstream), or an empty string.
        assert result.startswith("[error:") or result == ""

    def test_extract_text_pptx(self, tmp_path: Path):
        """Test extracting text from a .pptx file."""
        from pptx import Presentation

        pptx_file = tmp_path / "test.pptx"
        prs = Presentation()

        # Slide 1
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        for shape in slide1.shapes:
            if hasattr(shape, "text"):
                shape.text = "First Slide Title"

        # Slide 2
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        left = top = width = height = 1000000
        textbox = slide2.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "Bullet point content"

        prs.save(pptx_file)

        result = extract_text(pptx_file)
        assert result is not None
        assert not result.startswith("[error:")
        # Text content may vary depending on PowerPoint layout defaults
        assert len(result) > 0

    def test_extract_text_pptx_table(self, tmp_path: Path):
        """Table cells should be extracted, not silently dropped."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_file = tmp_path / "table.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = slide.shapes.add_table(
            2, 2, Inches(1), Inches(1), Inches(4), Inches(1)
        ).table
        table.cell(0, 0).text = "Header A"
        table.cell(0, 1).text = "Header B"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "Bob"
        prs.save(pptx_file)

        result = extract_text(pptx_file)
        assert result is not None
        assert "Header A" in result
        assert "Header B" in result
        assert "Alice" in result
        assert "Bob" in result

    def test_extract_text_pptx_grouped_shapes(self, tmp_path: Path):
        """Text inside grouped shapes must be extracted recursively."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_file = tmp_path / "grouped.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        group = slide.shapes.add_group_shape()
        inner = group.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        inner.text_frame.text = "Inside group"
        prs.save(pptx_file)

        result = extract_text(pptx_file)
        assert result is not None
        assert "Inside group" in result

    def test_extract_text_pdf_not_found(self, tmp_path: Path):
        """Test that missing PDF files return error string."""
        missing_pdf = tmp_path / "nonexistent.pdf"

        result = extract_text(missing_pdf)
        assert result is not None
        assert "[error: file not found:" in result

    def test_extract_text_image_files(self, tmp_path: Path):
        """Test that image files return placeholder text."""
        # Create a minimal PNG file (1x1 pixel)
        png_file = tmp_path / "test.png"
        # Minimal valid PNG: 8-byte signature + IHDR + IDAT + IEND
        png_data = (
            b"\x89PNG\r\n\x1a\n"
            b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01"
            b"\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        png_file.write_bytes(png_data)

        result = extract_text(png_file)
        assert result is not None
        assert "[image:" in result
        assert "test.png" in result


class TestIsTextExtension:
    """Test the _is_text_extension helper."""

    def test_text_extensions_return_true(self):
        """Test that known text extensions return True."""
        assert _is_text_extension(".txt") is True
        assert _is_text_extension(".md") is True
        assert _is_text_extension(".csv") is True
        assert _is_text_extension(".json") is True
        assert _is_text_extension(".yaml") is True
        assert _is_text_extension(".yml") is True
        assert _is_text_extension(".xml") is True
        assert _is_text_extension(".html") is True
        assert _is_text_extension(".htm") is True

    def test_non_text_extensions_return_false(self):
        """Test that non-text extensions return False."""
        assert _is_text_extension(".pdf") is False
        assert _is_text_extension(".docx") is False
        assert _is_text_extension(".xlsx") is False
        assert _is_text_extension(".pptx") is False
        assert _is_text_extension(".png") is False
        assert _is_text_extension(".xyz") is False

    def test_case_sensitivity(self):
        """Test that _is_text_extension requires lowercase extension.

        Note: The main extract_text function handles case-insensitivity by
        converting extensions to lowercase before calling _is_text_extension.
        """
        # _is_text_extension itself is case-sensitive (lowercase only)
        assert _is_text_extension(".txt") is True
        assert _is_text_extension(".TXT") is False
        assert _is_text_extension(".pdf") is False


def test_extract_documents_honours_an_explicit_size_limit(tmp_path):
    from durin.utils.document import extract_documents

    big = tmp_path / "big.txt"
    big.write_text("x" * 5000)

    text, images = extract_documents("hello", [str(big)], max_file_size=1000)
    assert "x" * 100 not in text
    assert images == []


def test_extract_documents_reports_a_skipped_oversized_file(tmp_path):
    # A silently dropped attachment reads to the user as durin ignoring them.
    from durin.utils.document import extract_documents

    big = tmp_path / "huge.pdf"
    big.write_bytes(b"%PDF-1.4" + b"0" * 5000)

    text, _ = extract_documents("hello", [str(big)], max_file_size=1000)
    assert "huge.pdf" in text
    assert "too large" in text.lower()


def test_configured_max_file_size_reads_the_configured_value(tmp_path, monkeypatch):
    # A non-default value, well clear of both the schema default (50) and the
    # module's fallback constant (also 50) — the two cannot be told apart by
    # a value that happens to equal either, so this pins the real read.
    import json

    from durin.utils.document import _configured_max_file_size

    (tmp_path / "config.json").write_text(
        json.dumps({"documents": {"maxFileSizeMb": 3}})
    )
    monkeypatch.setenv("DURIN_HOME", str(tmp_path))

    assert _configured_max_file_size() == 3 * 1024 * 1024


def test_configured_max_text_chars_reads_the_configured_value(tmp_path, monkeypatch):
    import json

    from durin.utils.document import _configured_max_text_chars

    (tmp_path / "config.json").write_text(
        json.dumps({"documents": {"maxTextChars": 1500}})
    )
    monkeypatch.setenv("DURIN_HOME", str(tmp_path))

    assert _configured_max_text_chars() == 1500


def test_extract_documents_reports_an_over_budget_scan_instead_of_dropping_it(
    tmp_path, monkeypatch
):
    # The flagship OCR case: a scanned book attached in chat with OCR on.
    # convert_file_to_markdown raises NeedsOcrJob because its page count
    # blows the inline budget; that used to vanish as an [error:] string the
    # drop gate discarded. The model must see the filename, the on-disk
    # path, and the "Ingest" instruction instead.
    import json

    from durin.utils.document import extract_documents
    from tests.tools.test_read_enhancements import _write_text_pdf

    monkeypatch.setenv("DURIN_HOME", str(tmp_path))
    (tmp_path / "config.json").write_text(
        json.dumps({"documents": {"ocr": {"enabled": True, "inline_max_pages": 2}}})
    )
    # The [ocr] extra is not part of CI's install set, so engine_available()
    # would read False there and this test would take the engine-missing
    # coverage-note branch instead of raising NeedsOcrJob — never exercising
    # the fix under test. Force it available so the test is deterministic
    # regardless of what is or isn't installed on the machine running it.
    monkeypatch.setattr("durin.memory.doc_convert.engine_available", lambda: True)

    pdf = tmp_path / "book.pdf"
    _write_text_pdf(pdf, [""] * 8)

    text, images = extract_documents("please remember this book", [str(pdf)])

    assert "book.pdf" in text
    assert str(pdf) in text
    assert "Ingest" in text
    assert images == []


def test_extract_documents_reports_a_blank_scan_instead_of_a_silent_empty_file(
    tmp_path, monkeypatch
):
    # A scanned PDF whose pages ALL transcribe blank must not vanish the way
    # an unraised empty string used to -- it has to name the file and say
    # why, like every other extraction failure this gate already covers.
    import json

    from durin.utils.document import extract_documents
    from tests.tools.test_read_enhancements import _write_text_pdf

    monkeypatch.setenv("DURIN_HOME", str(tmp_path))
    (tmp_path / "config.json").write_text(
        json.dumps({"documents": {"ocr": {"enabled": True, "inline_max_pages": 5}}})
    )
    # engine_available() would read False in CI (no [ocr] extra installed
    # there), taking the engine-missing coverage-note branch instead of ever
    # reaching the blank-after-OCR raise this test is about.
    monkeypatch.setattr("durin.memory.doc_convert.engine_available", lambda: True)
    monkeypatch.setattr(
        "durin.memory.doc_convert.transcribe_pages_detached",
        lambda path, pages: {p: "" for p in pages},
    )

    pdf = tmp_path / "blank_scan.pdf"
    _write_text_pdf(pdf, ["", ""])

    text, images = extract_documents("please remember this", [str(pdf)])

    assert "blank_scan.pdf" in text
    assert "could not be read inline" in text
    assert images == []


def test_extract_documents_reports_a_corrupt_file_instead_of_dropping_it(tmp_path):
    # The [error:] gate predates the OCR branch — it silently dropped
    # corrupt files too. The fix covers the whole defect class, not just
    # NeedsOcrJob.
    from durin.utils.document import extract_documents

    corrupt = tmp_path / "broken.pdf"
    corrupt.write_bytes(b"not a real pdf, just garbage bytes" * 5)

    text, images = extract_documents("hello", [str(corrupt)])

    assert "broken.pdf" in text
    assert "could not be read inline" in text
    assert images == []


def test_extract_documents_still_inlines_an_ordinary_pdf(tmp_path):
    # Regression: a normal text PDF must keep inlining cleanly with no
    # error line — the fix only changes what happens to [error:] strings.
    from durin.utils.document import extract_documents
    from tests.tools.test_read_enhancements import _write_text_pdf

    pdf = tmp_path / "report.pdf"
    _write_text_pdf(pdf, ["Quarterly revenue is $5M"])

    text, images = extract_documents("summarize", [str(pdf)])

    assert "Quarterly revenue is $5M" in text
    assert "[error:" not in text
    assert "could not be read inline" not in text
    assert images == []
